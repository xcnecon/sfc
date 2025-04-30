import os
from pptx import Presentation
from pptx.util import Inches, Pt

# Set up directories and parameters
output_dir = "output"
prefixes = ["wealth", "gpd", "ratio", "region_gdp", "di", "consumption", "transfer"]
subs = range(1, 8)  # 1 to 7

# Scenario titles as in baseline.prg
scenario_titles = [
    "Shock 1: increase transfer",
    "Shock 2: reduce Import",
    "Shock 3: reduce transfer",
    "Shock 4: increase government expenditure in west",
    "Shock 5: increase government expenditure in east",
    "Shock 6: increase loan rate",
    "Shock 7: increase tax",
]

prs = Presentation()
title_slide_layout = prs.slide_layouts[5]  # Blank slide

# Image size and grid parameters
img_width = Inches(2.2)
img_height = Inches(2.2)
cols = 3
x_margin = Inches(0.5)
y_margin = Inches(1.0)
x_spacing = Inches(0.2)
y_spacing = Inches(0.2)

for idx, sub in enumerate(subs):
    slide = prs.slides.add_slide(title_slide_layout)
    # Add a title from scenario_titles
    left = Inches(0.5)
    top = Inches(0.2)
    title_shape = slide.shapes.add_textbox(left, top, Inches(8), Inches(0.5))
    title_frame = title_shape.text_frame
    title_frame.text = scenario_titles[idx]
    title_frame.paragraphs[0].font.size = Pt(24)
    
    # Place images in a grid
    for i, prefix in enumerate(prefixes):
        fname = f"{prefix}_{sub}.emf"
        fpath = os.path.join(output_dir, fname)
        if os.path.exists(fpath):
            row = i // cols
            col = i % cols
            img_left = x_margin + col * (img_width + x_spacing)
            img_top = y_margin + row * (img_height + y_spacing)
            slide.shapes.add_picture(fpath, img_left, img_top, img_width, img_height)
            # Add a caption below each image
            caption_top = img_top + img_height
            caption_shape = slide.shapes.add_textbox(img_left, caption_top, img_width, Inches(0.3))
            caption_frame = caption_shape.text_frame
            caption_frame.text = prefix.capitalize()
            caption_frame.paragraphs[0].font.size = Pt(12)

prs.save("simulation_results.pptx")
print("PPT file created: simulation_results.pptx")